import os, json, re, time, requests, streamlit as st
from pptx import Presentation
from dotenv import load_dotenv
from io import BytesIO

load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "")
OPENROUTER_URL     = "https://openrouter.ai/api/v1/chat/completions"
MODEL              = "anthropic/claude-haiku-4-5"

DIFF_HINTS = {
    "Simple":  "Basic recall — key terms, definitions straight from slides.",
    "Medium":  "Balanced mix of recall + scenario-based reasoning.",
    "Complex": "Analytical edge-cases, trade-offs & deep application.",
}
DIFF_PROMPTS = {
    "Simple":  "Generate straightforward factual recall questions focusing on definitions, key terms, and basic concepts directly stated in the slides.",
    "Medium":  "Generate balanced questions mixing factual recall and scenario-based reasoning.",
    "Complex": "Generate challenging analytical questions requiring deep understanding, edge cases, trade-offs, and novel application.",
}

st.set_page_config(
    page_title="AI Quiz Generator",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ═══════════════════════════════════════════════════════════════════════════════
#  GLOBAL CSS — PREMIUM DARK THEME
# ═══════════════════════════════════════════════════════════════════════════════
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');

/* ── RESET & BASE ── */
*, *::before, *::after { box-sizing: border-box; }

html, body {
    font-family: 'Inter', sans-serif;
    background: #080c14;
    color: #f1f5f9;
    min-height: 100vh;
    overflow: auto;
}

[data-testid="stAppViewContainer"] {
    background: #080c14;
    font-family: 'Inter', sans-serif;
    min-height: 100vh;
}

[data-testid="stHeader"]  { display: none !important; }
[data-testid="stSidebar"] { display: none !important; }
#MainMenu, footer         { visibility: hidden; }

[data-testid="stMain"] {
    padding: 0 !important;
    min-height: 100vh;
}
[data-testid="stMain"] > div {
    padding: 0 !important;
    min-height: 100vh;
}
[data-testid="block-container"] {
    padding: 0 !important;
    max-width: 100% !important;
    min-height: 100vh;
}

/* let Streamlit columns breathe — never clip */
[data-testid="column"],
[data-testid="stVerticalBlock"],
[data-testid="stVerticalBlockBorderWrapper"],
[data-testid="stHorizontalBlock"] {
    overflow: visible !important;
    height: auto !important;
    min-height: 0 !important;
}

/* ── LAYOUT SHELL ── */
.shell {
    display: flex;
    min-height: 100vh;
    background: #080c14;
}

/* ── LEFT NAV DOCK ── */
.nav-dock {
    width: 72px;
    min-height: 100vh;
    background: rgba(255,255,255,0.03);
    border-right: 1px solid rgba(255,255,255,0.06);
    display: flex;
    flex-direction: column;
    align-items: center;
    padding: 20px 0;
    gap: 8px;
    flex-shrink: 0;
    backdrop-filter: blur(20px);
    position: sticky;
    top: 0;
}
.nav-logo {
    width: 40px; height: 40px;
    background: linear-gradient(135deg, #6366f1, #8b5cf6);
    border-radius: 12px;
    display: flex; align-items: center; justify-content: center;
    font-size: 20px;
    box-shadow: 0 0 20px rgba(99,102,241,0.5);
    margin-bottom: 16px;
}
.nav-sep { width: 32px; height: 1px; background: rgba(255,255,255,0.07); margin: 8px 0; }
.nav-item {
    width: 44px; height: 44px;
    border-radius: 12px;
    display: flex; align-items: center; justify-content: center;
    font-size: 19px; cursor: pointer;
    transition: all .2s; color: rgba(255,255,255,0.35);
    position: relative;
}
.nav-item:hover { background: rgba(255,255,255,0.07); color: rgba(255,255,255,0.8); }
.nav-item.active { background: rgba(99,102,241,0.2); color: #a5b4fc; }
.nav-item.active::before {
    content: '';
    position: absolute; left: -12px;
    width: 3px; height: 24px;
    background: linear-gradient(180deg,#6366f1,#8b5cf6);
    border-radius: 0 4px 4px 0;
}

/* ── MAIN PANEL ── */
.main-panel {
    flex: 1;
    display: flex;
    flex-direction: column;
    min-height: 100vh;
}

/* ── TOP HEADER ── */
.top-header {
    height: 64px;
    background: rgba(255,255,255,0.02);
    border-bottom: 1px solid rgba(255,255,255,0.06);
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 0 36px;
    flex-shrink: 0;
}
.header-title {
    font-size: 15px;
    font-weight: 700;
    color: #f1f5f9;
    display: flex; align-items: center; gap: 10px;
}
.header-title span { color: rgba(255,255,255,0.35); font-weight: 400; }

/* STEP PROGRESS BAR */
.step-track {
    display: flex;
    align-items: center;
    gap: 0;
}
.step-node {
    display: flex; align-items: center; gap: 8px;
}
.step-circle {
    width: 28px; height: 28px;
    border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-size: 11px; font-weight: 800;
    border: 1.5px solid rgba(255,255,255,0.12);
    color: rgba(255,255,255,0.3);
    background: transparent;
    transition: all .3s;
}
.step-circle.done   { background: #6366f1; border-color: #6366f1; color: #fff; box-shadow: 0 0 12px rgba(99,102,241,0.5); }
.step-circle.active { background: transparent; border-color: #6366f1; color: #818cf8; box-shadow: 0 0 0 4px rgba(99,102,241,0.15); }
.step-label { font-size: 11px; font-weight: 600; color: rgba(255,255,255,0.3); white-space: nowrap; }
.step-label.active { color: #818cf8; }
.step-line { width: 40px; height: 1px; background: rgba(255,255,255,0.1); margin: 0 8px; }
.step-line.done { background: #6366f1; }

.header-right { display: flex; align-items: center; gap: 12px; }
.hdr-avatar {
    width: 32px; height: 32px; border-radius: 50%;
    background: linear-gradient(135deg,#6366f1,#8b5cf6);
    display: flex; align-items: center; justify-content: center;
    font-size: 14px;
}

/* ── CONTENT AREA ── */
.content-area {
    flex: 1;
    display: flex;
    align-items: flex-start;
    justify-content: center;
    padding: 24px 36px 40px;
}

/* ── GLASS PANEL ── */
.glass-panel {
    background: rgba(255,255,255,0.03);
    backdrop-filter: blur(24px);
    -webkit-backdrop-filter: blur(24px);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 24px;
    width: 100%;
    max-width: 660px;
    padding: 32px 40px;
    box-shadow:
        0 0 0 1px rgba(255,255,255,0.04) inset,
        0 24px 64px rgba(0,0,0,0.5),
        0 4px 16px rgba(0,0,0,0.3);
    position: relative;
    overflow: visible;
}
.glass-panel::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 1px;
    background: linear-gradient(90deg, transparent, rgba(99,102,241,0.5), rgba(139,92,246,0.4), transparent);
}

/* ── SCREEN HEADING ── */
.screen-head { margin-bottom: 20px; }
.screen-head .eyebrow {
    font-size: 10px; font-weight: 700; letter-spacing: .14em;
    text-transform: uppercase;
    background: linear-gradient(90deg,#818cf8,#a78bfa);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    margin-bottom: 10px; display: block;
}
.screen-head h2 {
    font-size: 26px; font-weight: 800; color: #f1f5f9;
    letter-spacing: -.5px; margin-bottom: 7px; line-height: 1.2;
}
.screen-head p { font-size: 14px; color: rgba(255,255,255,.4); line-height: 1.5; }

/* ── UPLOAD ZONE ── */
@keyframes borderGlow {
    0%,100% { opacity:.5; }
    50%      { opacity:1; }
}
@keyframes floatIcon {
    0%,100% { transform: translateY(0); }
    50%      { transform: translateY(-8px); }
}
@keyframes scanLine {
    0%   { top: 0%; opacity:1; }
    100% { top: 100%; opacity:0; }
}
@keyframes pulseRing {
    0%   { transform:scale(1); opacity:.6; }
    100% { transform:scale(1.8); opacity:0; }
}
@keyframes spin {
    to { transform: rotate(360deg); }
}
@keyframes shimmer {
    0%   { background-position: -200% center; }
    100% { background-position:  200% center; }
}

.upload-zone {
    border-radius: 18px;
    padding: 30px 28px;
    text-align: center;
    cursor: pointer;
    position: relative;
    overflow: hidden;
    transition: all .3s;
    background: rgba(99,102,241,0.04);
    border: 1.5px dashed rgba(99,102,241,0.35);
    margin-bottom: 14px;
}
.upload-zone::before {
    content: '';
    position: absolute; inset: 0;
    background: radial-gradient(circle at 50% 0%, rgba(99,102,241,0.1), transparent 60%);
    animation: borderGlow 3s ease-in-out infinite;
    pointer-events: none;
}
.upload-zone:hover {
    border-color: rgba(99,102,241,0.7);
    background: rgba(99,102,241,0.08);
    box-shadow: 0 0 40px rgba(99,102,241,0.15);
}
.upload-icon-wrap {
    display: inline-flex; align-items: center; justify-content: center;
    width: 72px; height: 72px; border-radius: 20px;
    background: rgba(99,102,241,0.12);
    border: 1px solid rgba(99,102,241,0.25);
    margin: 0 auto 16px;
    font-size: 32px;
    animation: floatIcon 3s ease-in-out infinite;
}
.upload-zone h3 { font-size: 18px; font-weight: 700; color: #f1f5f9; margin-bottom: 7px; }
.upload-zone p  { font-size: 13px; color: rgba(255,255,255,.35); line-height: 1.6; }
.upload-zone .constraint {
    display: inline-flex; align-items: center; gap: 6px;
    margin-top: 14px;
    background: rgba(255,255,255,0.05);
    border: 1px solid rgba(255,255,255,0.08);
    border-radius: 20px; padding: 5px 14px;
    font-size: 11px; font-weight: 600; color: rgba(255,255,255,.35);
    letter-spacing: .04em;
}

/* Browse pill button */
.browse-pill {
    display: flex; align-items: center; justify-content: center; gap: 8px;
    width: 100%; padding: 13px 20px;
    border-radius: 14px;
    background: rgba(255,255,255,0.04);
    border: 1px solid rgba(255,255,255,0.09);
    color: rgba(255,255,255,.65);
    font-size: 13px; font-weight: 600;
    cursor: pointer; transition: all .2s;
    margin-bottom: 16px;
}
.browse-pill:hover { border-color: rgba(99,102,241,0.5); color: #a5b4fc; background: rgba(99,102,241,0.08); }

/* Parsed file card */
.file-card {
    display: flex; align-items: center; gap: 14px;
    background: rgba(16,185,129,0.08);
    border: 1px solid rgba(16,185,129,0.22);
    border-radius: 14px; padding: 14px 18px;
    margin: 12px 0;
    box-shadow: 0 0 20px rgba(16,185,129,0.07);
}
.file-card .fc-icon { font-size: 22px; }
.file-card .fc-info { flex: 1; }
.file-card .fc-info strong { font-size: 13px; font-weight: 700; color: #f1f5f9; display: block; margin-bottom: 2px; }
.file-card .fc-info span   { font-size: 11px; color: rgba(255,255,255,.4); }
.fc-badge { background: linear-gradient(135deg,#059669,#10b981); color: #fff; font-size: 10px; font-weight: 800; padding: 4px 12px; border-radius: 20px; letter-spacing: .04em; box-shadow: 0 0 12px rgba(16,185,129,0.4); }

/* ── AI PROCESSING SCREEN ── */
.processing-wrap { text-align: center; padding: 20px 0; }
.pulse-ring-wrap {
    position: relative; width: 100px; height: 100px;
    margin: 0 auto 28px;
}
.pulse-ring {
    position: absolute; inset: 0; border-radius: 50%;
    border: 2px solid rgba(99,102,241,0.6);
    animation: pulseRing 1.8s ease-out infinite;
}
.pulse-ring:nth-child(2) { animation-delay: .6s; }
.pulse-ring:nth-child(3) { animation-delay: 1.2s; }
.pulse-core {
    position: absolute; inset: 20px;
    border-radius: 50%;
    background: linear-gradient(135deg,#6366f1,#8b5cf6);
    display: flex; align-items: center; justify-content: center;
    font-size: 22px;
    box-shadow: 0 0 30px rgba(99,102,241,0.6);
    animation: spin 3s linear infinite;
}
.processing-title { font-size: 22px; font-weight: 800; color: #f1f5f9; margin-bottom: 8px; }
.processing-sub   { font-size: 13px; color: rgba(255,255,255,.4); margin-bottom: 28px; }
.proc-steps { display: flex; flex-direction: column; gap: 10px; }
.proc-step {
    display: flex; align-items: center; gap: 12px;
    background: rgba(255,255,255,0.04);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 12px; padding: 12px 16px;
    font-size: 13px; font-weight: 500; color: rgba(255,255,255,.5);
    transition: all .4s;
}
.proc-step.active {
    background: rgba(99,102,241,0.12);
    border-color: rgba(99,102,241,0.3);
    color: #c7d2fe;
}
.proc-step.done {
    background: rgba(16,185,129,0.08);
    border-color: rgba(16,185,129,0.25);
    color: #6ee7b7;
}
.step-dot {
    width: 8px; height: 8px; border-radius: 50%;
    background: rgba(255,255,255,0.15); flex-shrink: 0;
}
.proc-step.active .step-dot { background: #6366f1; box-shadow: 0 0 8px #6366f1; animation: pulseRing .8s ease-out infinite; }
.proc-step.done   .step-dot { background: #10b981; }

/* ── SOURCE PILL ── */
.src-pill {
    display: flex; align-items: center; gap: 10px;
    background: rgba(99,102,241,0.08);
    border: 1px solid rgba(99,102,241,0.2);
    border-radius: 12px; padding: 11px 16px;
    font-size: 13px; font-weight: 600; color: #a5b4fc;
    margin-bottom: 24px;
}
.src-pill em { color: rgba(255,255,255,.35); font-style: normal; font-weight: 400; }

/* ── SECTION LABEL ── */
.slbl { font-size: 10px; font-weight: 700; letter-spacing: .12em; text-transform: uppercase; color: rgba(255,255,255,.3); margin-bottom: 12px; display: block; }

/* ── DIFFICULTY GRID ── */
.diff-grid { display: flex; gap: 10px; margin-bottom: 12px; }
.diff-tile {
    flex: 1; padding: 16px 10px; border-radius: 14px;
    border: 1.5px solid rgba(255,255,255,0.07);
    background: rgba(255,255,255,0.03);
    text-align: center; cursor: pointer; transition: all .2s;
}
.diff-tile:hover { border-color: rgba(99,102,241,0.4); background: rgba(99,102,241,0.06); }
.diff-tile.on {
    border-color: #6366f1;
    background: rgba(99,102,241,0.12);
    box-shadow: 0 0 0 3px rgba(99,102,241,0.15);
}
.diff-tile .de { font-size: 22px; margin-bottom: 6px; display: block; }
.diff-tile .dn { font-size: 12px; font-weight: 700; color: rgba(255,255,255,.5); }
.diff-tile.on .dn { color: #a5b4fc; }
.hint-card {
    background: rgba(6,182,212,0.06); border: 1px solid rgba(6,182,212,0.15);
    border-radius: 10px; padding: 10px 14px; font-size: 12px; color: #67e8f9;
    display: flex; gap: 8px; margin-bottom: 22px;
}

/* ── QUIZ ONE-QUESTION LAYOUT ── */
.quiz-meta {
    display: flex; align-items: center; justify-content: space-between;
    margin-bottom: 10px;
}
.qm-counter {
    font-size: 11px; font-weight: 700; letter-spacing: .1em;
    background: linear-gradient(90deg,#818cf8,#a78bfa);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
}
.qm-timer {
    background: rgba(251,146,60,0.1); border: 1px solid rgba(251,146,60,0.25);
    color: #fb923c; font-size: 13px; font-weight: 700;
    padding: 5px 13px; border-radius: 10px;
    display: flex; align-items: center; gap: 5px;
}

.pbar { background: rgba(255,255,255,0.07); border-radius: 6px; height: 5px; margin-bottom: 24px; overflow: hidden; }
.pbar-f {
    height: 100%; border-radius: 6px;
    background: linear-gradient(90deg,#6366f1,#8b5cf6);
    box-shadow: 0 0 10px rgba(99,102,241,0.5);
    transition: width .5s ease;
}

.q-topic { display: inline-flex; background: rgba(99,102,241,0.12); border: 1px solid rgba(99,102,241,0.25); color: #a5b4fc; font-size: 10px; font-weight: 700; letter-spacing: .06em; padding: 4px 10px; border-radius: 6px; margin-bottom: 12px; }
.q-text { font-size: 17px; font-weight: 700; color: #f1f5f9; line-height: 1.55; margin-bottom: 22px; letter-spacing: -.2px; }

/* Options as big clickable cards */
.opt-card {
    display: flex; align-items: center; gap: 14px;
    width: 100%; padding: 15px 18px;
    border-radius: 14px; cursor: pointer;
    margin-bottom: 10px; transition: all .18s;
    border: 1.5px solid rgba(255,255,255,0.08);
    background: rgba(255,255,255,0.03);
    color: rgba(255,255,255,.75); font-size: 14px; font-weight: 500;
}
.opt-card:hover { border-color: rgba(99,102,241,0.45); background: rgba(99,102,241,0.08); color: #c7d2fe; }
.opt-card.picked {
    border-color: #6366f1 !important;
    background: rgba(99,102,241,0.15) !important;
    color: #e0e7ff !important; font-weight: 600 !important;
    box-shadow: 0 0 0 3px rgba(99,102,241,0.2), 0 0 20px rgba(99,102,241,0.1) !important;
}
.opt-badge {
    width: 30px; height: 30px; border-radius: 8px; flex-shrink: 0;
    display: flex; align-items: center; justify-content: center;
    font-size: 12px; font-weight: 800;
    background: rgba(255,255,255,0.07); color: rgba(255,255,255,.4);
    border: 1px solid rgba(255,255,255,0.1);
    transition: all .18s;
}
.opt-card.picked .opt-badge { background: linear-gradient(135deg,#6366f1,#8b5cf6); color: #fff; border-color: transparent; box-shadow: 0 0 12px rgba(99,102,241,0.5); }

/* ── RESULTS ── */
.result-banner {
    background: linear-gradient(135deg,#1e1b4b,#312e81,#1e40af);
    border-radius: 20px; padding: 36px 28px; text-align: center;
    margin-bottom: 20px; position: relative; overflow: hidden;
    box-shadow: 0 8px 40px rgba(99,102,241,0.3);
}
.result-banner::before { content:''; position:absolute; top:-60px; right:-60px; width:200px; height:200px; border-radius:50%; background:rgba(255,255,255,0.05); }
.result-banner::after  { content:''; position:absolute; bottom:-80px; left:-40px; width:240px; height:240px; border-radius:50%; background:rgba(255,255,255,0.03); }
.rb-score {
    font-size: 72px; font-weight: 900; line-height: 1;
    background: linear-gradient(135deg,#a5b4fc,#c4b5fd,#93c5fd);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    position: relative; z-index: 1;
    filter: drop-shadow(0 0 20px rgba(165,180,252,0.4));
}
.rb-pct  { font-size: 20px; color: rgba(255,255,255,.6); font-weight: 600; margin-bottom: 6px; position: relative; z-index: 1; }
.rb-msg  { font-size: 16px; font-weight: 700; color: #fff; margin-bottom: 18px; position: relative; z-index: 1; }
.rb-chips { display: flex; justify-content: center; gap: 10px; flex-wrap: wrap; position: relative; z-index: 1; }
.rbc { padding: 6px 16px; border-radius: 30px; font-size: 12px; font-weight: 700; }
.rbc-g { background: rgba(16,185,129,.2); color: #6ee7b7; border: 1px solid rgba(16,185,129,.35); }
.rbc-r { background: rgba(239,68,68,.18); color: #fca5a5; border: 1px solid rgba(239,68,68,.3); }
.rbc-b { background: rgba(255,255,255,.1); color: rgba(255,255,255,.75); border: 1px solid rgba(255,255,255,.18); }

.rev-section { font-size: 10px; font-weight: 700; letter-spacing: .12em; color: rgba(255,255,255,.25); text-transform: uppercase; margin: 18px 0 10px; }
.rev-scroll  { max-height: 280px; overflow-y: auto; padding-right: 4px; }
.rev-scroll::-webkit-scrollbar { width: 3px; }
.rev-scroll::-webkit-scrollbar-track { background: rgba(255,255,255,0.04); border-radius: 4px; }
.rev-scroll::-webkit-scrollbar-thumb { background: rgba(255,255,255,0.12); border-radius: 4px; }

.rv { border-radius: 12px; padding: 13px 16px; margin-bottom: 8px; }
.rv.ok  { background: rgba(16,185,129,0.07); border: 1px solid rgba(16,185,129,0.2); }
.rv.bad { background: rgba(239,68,68,0.07);  border: 1px solid rgba(239,68,68,0.18); }
.rv-hd  { display: flex; align-items: flex-start; gap: 9px; margin-bottom: 4px; }
.rv-tl  { font-size: 12px; font-weight: 700; color: #f1f5f9; }
.rv-q   { font-size: 11px; color: rgba(255,255,255,.38); margin: 3px 0 4px 22px; line-height: 1.4; }
.rv-ok  { font-size: 11px; font-weight: 600; color: #34d399; margin-left: 22px; }
.rv-bad { font-size: 11px; font-weight: 600; color: #f87171; margin-left: 22px; }
.rv-ai  { background: rgba(99,102,241,0.08); border: 1px solid rgba(99,102,241,0.18); border-radius: 8px; padding: 8px 12px; font-size: 11px; color: #c7d2fe; margin: 6px 0 0; display: flex; gap: 7px; line-height: 1.4; }

/* ══════════════════════════════════════════
   STREAMLIT WIDGET OVERRIDES
══════════════════════════════════════════ */
div[data-testid="stFileUploaderDropzone"] {
    background: rgba(99,102,241,0.05) !important;
    border: 1.5px dashed rgba(99,102,241,0.35) !important;
    border-radius: 16px !important;
}
div[data-testid="stFileUploaderDropzone"] p     { color: rgba(255,255,255,.4) !important; }
div[data-testid="stFileUploaderDropzone"] small { color: rgba(255,255,255,.25) !important; }
div[data-testid="stFileUploaderDropzone"] svg   { fill: rgba(99,102,241,.5) !important; }

div[data-testid="stSlider"] label { color: rgba(255,255,255,.75) !important; font-weight: 600 !important; font-size: 13px !important; }
div[data-testid="stSlider"] p { color: rgba(255,255,255,.35) !important; }
[data-testid="stSlider"] [role="slider"] { background: #6366f1 !important; box-shadow: 0 0 0 4px rgba(99,102,241,.25) !important; }
[data-testid="stSlider"] [data-baseweb="slider"] [role="progressbar"] { background: linear-gradient(90deg,#6366f1,#8b5cf6) !important; }

/* ── ALL BUTTONS BASE ── */
.stButton > button {
    border-radius: 12px !important;
    font-weight: 600 !important;
    font-size: 14px !important;
    padding: 13px 20px !important;
    transition: all .2s !important;
    width: 100% !important;
    font-family: 'Inter', sans-serif !important;
    text-align: left !important;
}

/* ── PRIMARY (action buttons + selected options) ── */
.stButton > button[kind="primary"] {
    background: linear-gradient(135deg,#6366f1,#8b5cf6) !important;
    color: #fff !important;
    border: none !important;
    box-shadow: 0 4px 20px rgba(99,102,241,0.45) !important;
}
.stButton > button[kind="primary"]:hover {
    box-shadow: 0 6px 30px rgba(99,102,241,0.65) !important;
    transform: translateY(-2px) !important;
}

/* ── SECONDARY (unselected options + nav buttons) ── */
.stButton > button[kind="secondary"] {
    background: rgba(255,255,255,0.03) !important;
    border: 1.5px solid rgba(255,255,255,0.09) !important;
    color: rgba(255,255,255,.7) !important;
}
.stButton > button[kind="secondary"]:hover {
    background: rgba(99,102,241,0.09) !important;
    border-color: rgba(99,102,241,0.45) !important;
    color: #c7d2fe !important;
    transform: translateY(-1px) !important;
}

.stButton > button:disabled { opacity: .3 !important; transform: none !important; box-shadow: none !important; }

/* ── NAV buttons (Back / Next) — centred text ── */
.stButton > button[data-testid*="Back"],
.stButton > button[data-testid*="Next"],
.stButton > button[data-testid*="Submit"],
.stButton > button[data-testid*="Retake"],
.stButton > button[data-testid*="Upload"] {
    text-align: center !important;
}

[data-testid="stSpinner"] > div { border-top-color: #6366f1 !important; }
.stCaption, [data-testid="stCaptionContainer"] p { color: rgba(255,255,255,.28) !important; font-size: 11px !important; }
[data-testid="stAlert"] { border-radius: 12px !important; }
.element-container { margin: 0 !important; }

/* ── NAV DOCK BUTTONS ── */
div[data-testid="stVerticalBlock"] > div[data-testid="stVerticalBlockBorderWrapper"] > div > div > div.stButton > button,
section[data-testid="stSidebar"] .stButton > button {
    width: 44px !important; height: 44px !important;
    min-height: 44px !important;
    padding: 0 !important;
    border-radius: 12px !important;
    font-size: 20px !important;
    display: flex !important; align-items: center !important; justify-content: center !important;
    text-align: center !important;
    margin: 2px auto !important;
}

/* nav column specific overrides */
.nav-btn-col .stButton > button {
    width: 44px !important; height: 44px !important;
    min-height: 44px !important;
    padding: 0 !important;
    border-radius: 12px !important;
    font-size: 20px !important;
    text-align: center !important;
    line-height: 1 !important;
    margin: 0 auto !important;
}
.nav-btn-col .stButton > button[kind="primary"] {
    background: rgba(99,102,241,0.2) !important;
    border: 1px solid rgba(99,102,241,0.4) !important;
    color: #a5b4fc !important;
    box-shadow: none !important;
    transform: none !important;
}
.nav-btn-col .stButton > button[kind="secondary"] {
    background: transparent !important;
    border: 1px solid transparent !important;
    color: rgba(255,255,255,0.3) !important;
}
.nav-btn-col .stButton > button[kind="secondary"]:hover {
    background: rgba(255,255,255,0.07) !important;
    border-color: rgba(255,255,255,0.1) !important;
    color: rgba(255,255,255,0.8) !important;
    transform: none !important;
}
.nav-btn-col [data-testid="stVerticalBlock"] { gap: 4px !important; }

/* info/stat cards for analytics/history/settings pages */
.info-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 12px; margin-bottom: 20px; }
.info-card {
    background: rgba(255,255,255,0.04); border: 1px solid rgba(255,255,255,0.08);
    border-radius: 16px; padding: 20px 18px;
}
.info-card .ic-label { font-size: 10px; font-weight: 700; letter-spacing: .1em; text-transform: uppercase; color: rgba(255,255,255,.3); margin-bottom: 8px; }
.info-card .ic-val   { font-size: 28px; font-weight: 800; color: #f1f5f9; }
.info-card .ic-sub   { font-size: 12px; color: rgba(255,255,255,.3); margin-top: 3px; }
.hist-row {
    display: flex; align-items: center; justify-content: space-between;
    background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.07);
    border-radius: 12px; padding: 14px 18px; margin-bottom: 8px;
}
.hist-row .hr-title { font-size: 13px; font-weight: 600; color: #f1f5f9; }
.hist-row .hr-sub   { font-size: 11px; color: rgba(255,255,255,.35); margin-top: 2px; }
.hist-score { font-size: 18px; font-weight: 800; color: #a5b4fc; }
.setting-row {
    display: flex; align-items: center; justify-content: space-between;
    background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.07);
    border-radius: 12px; padding: 16px 20px; margin-bottom: 8px;
}
.setting-row .sr-label { font-size: 13px; font-weight: 600; color: #f1f5f9; }
.setting-row .sr-sub   { font-size: 11px; color: rgba(255,255,255,.35); margin-top: 2px; }
.setting-row .sr-val   { font-size: 12px; font-weight: 700; color: #818cf8; background: rgba(99,102,241,0.12); border: 1px solid rgba(99,102,241,0.25); padding: 4px 12px; border-radius: 8px; }
</style>
""", unsafe_allow_html=True)

# ── SESSION STATE ─────────────────────────────────────────────────────────────
def init():
    for k, v in {
        "step":"upload","slides":[],"filename":"","slide_count":0,"word_count":0,
        "num_questions":10,"difficulty":"Medium","questions":[],"user_answers":{},
        "current_q":0,"start_time":None,"elapsed":0,"explanations":{},"rk":0,
        "processing":False,"proc_stage":0,"nav_page":"home",
        "quiz_history":[],
    }.items():
        if k not in st.session_state: st.session_state[k] = v
init()

# ── BACKEND ───────────────────────────────────────────────────────────────────
def extract_pptx(b):
    prs = Presentation(BytesIO(b)); slides, words = [], 0
    for i, slide in enumerate(prs.slides, 1):
        txt = [p.text.strip() for s in slide.shapes if s.has_text_frame for p in s.text_frame.paragraphs if p.text.strip()]
        c = " ".join(txt)
        if c: slides.append({"slide":i,"text":c}); words += len(c.split())
    return slides, words

def llm(prompt):
    r = requests.post(OPENROUTER_URL,
        headers={"Authorization":f"Bearer {OPENROUTER_API_KEY}","Content-Type":"application/json","HTTP-Referer":"http://localhost:8501","X-Title":"AI Quiz Generator"},
        json={"model":MODEL,"messages":[{"role":"user","content":prompt}],"temperature":.7,"max_tokens":4096},timeout=120)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"]

def gen_mcqs(slides, n, diff):
    txt = "\n\n".join(f"[Slide {s['slide']}]: {s['text']}" for s in slides)
    raw = llm(f"Expert quiz designer. Generate exactly {n} MCQs.\nDIFFICULTY: {diff} — {DIFF_PROMPTS[diff]}\nSLIDES:\n{txt}\nReturn ONLY valid JSON array, no markdown:\n[{{\"id\":1,\"question\":\"...\",\"options\":{{\"A\":\"...\",\"B\":\"...\",\"C\":\"...\",\"D\":\"...\"}},\"correct\":\"B\",\"topic\":\"label\"}}]")
    return json.loads(re.sub(r"```(?:json)?","",raw).strip().strip("`").strip())

def gen_exp(questions, answers):
    wrong = [q for q in questions if answers.get(str(q["id"])) and answers[str(q["id"])] != q["correct"]]
    if not wrong: return {}
    block = "\n\n".join(f"Q{w['id']}: {w['question']}\nStudent: {answers[str(w['id'])]} — {w['options'].get(answers[str(w['id'])],'')} \nCorrect: {w['correct']} — {w['options'][w['correct']]}" for w in wrong)
    raw = llm(f"1-2 sentences per wrong answer: why wrong + why correct is right.\n{block}\nReturn ONLY JSON: {{\"1\":\"...\"}}")
    return json.loads(re.sub(r"```(?:json)?","",raw).strip().strip("`").strip())

# ── NAV & STEP CONFIG ────────────────────────────────────────────────────────
NAV_PAGES = [("🏠","home"),("📊","analytics"),("🕘","history"),("⚙️","settings")]

step_order = ["upload","config","quiz","results"]
step_names = ["Upload","Configure","Quiz","Results"]
ci = step_order.index(st.session_state.step)

def steps_html():
    out = ""
    for idx,(_,name) in enumerate(zip(step_order,step_names)):
        dc = "done" if idx<ci else ("active" if idx==ci else "")
        lc = "active" if idx==ci else ""
        lbl = "✓" if idx<ci else str(idx+1)
        out += f'<div class="step-node"><div class="step-circle {dc}">{lbl}</div><span class="step-label {lc}">{name}</span></div>'
        if idx<3: out += f'<div class="step-line {"done" if idx<ci else ""}"></div>'
    return out

nav_col, main_col = st.columns([0.06, 0.94])

with nav_col:
    st.markdown('<div class="nav-btn-col">', unsafe_allow_html=True)
    st.markdown('<div style="text-align:center;padding:14px 0 10px"><div class="nav-logo" style="margin:0 auto">🎯</div></div>', unsafe_allow_html=True)
    st.markdown('<div style="width:32px;height:1px;background:rgba(255,255,255,0.07);margin:4px auto 8px"></div>', unsafe_allow_html=True)
    for icon, page in NAV_PAGES:
        is_active = st.session_state.nav_page == page
        if st.button(icon, key=f"nav_{page}", type="primary" if is_active else "secondary", help=page.capitalize()):
            st.session_state.nav_page = page
            if page == "home":
                pass  # stay on current quiz step
            st.rerun()
    st.markdown('</div>', unsafe_allow_html=True)

with main_col:
    nav = st.session_state.nav_page

    # TOP HEADER
    page_title = {"home": step_names[ci], "analytics": "Analytics", "history": "History", "settings": "Settings"}[nav]
    st.markdown(f"""
    <div class="top-header">
      <div class="header-title">AI Quiz Generator <span>/</span> <span style="color:rgba(255,255,255,.6)">{page_title}</span></div>
      <div class="step-track">{"" if nav != "home" else steps_html()}</div>
      <div class="header-right"><div class="hdr-avatar">👤</div></div>
    </div>
    """, unsafe_allow_html=True)

    # CONTENT AREA — centered glass panel
    _, center, _ = st.columns([1, 3, 1])
    with center:

        # ══════════════════════════════════════════════════════════════════════
        # ANALYTICS PAGE
        # ══════════════════════════════════════════════════════════════════════
        if nav == "analytics":
            hist = st.session_state.quiz_history
            total_quizzes = len(hist)
            avg_score = round(sum(h["pct"] for h in hist)/total_quizzes) if hist else 0
            best = max((h["pct"] for h in hist), default=0)
            total_q = sum(h["total"] for h in hist)

            st.markdown(f"""
            <div class="screen-head">
              <span class="eyebrow">Your Performance</span>
              <h2>Analytics Dashboard</h2>
              <p>Stats from all your quiz sessions this session.</p>
            </div>
            <div class="info-grid">
              <div class="info-card"><div class="ic-label">Quizzes Taken</div><div class="ic-val">{total_quizzes}</div><div class="ic-sub">this session</div></div>
              <div class="info-card"><div class="ic-label">Avg Score</div><div class="ic-val">{avg_score}%</div><div class="ic-sub">across all quizzes</div></div>
              <div class="info-card"><div class="ic-label">Best Score</div><div class="ic-val">{best}%</div><div class="ic-sub">personal best</div></div>
              <div class="info-card"><div class="ic-label">Total Questions</div><div class="ic-val">{total_q}</div><div class="ic-sub">answered</div></div>
            </div>
            """, unsafe_allow_html=True)
            if not hist:
                st.markdown('<div style="color:rgba(255,255,255,.3);font-size:13px;text-align:center;padding:20px">Complete a quiz to see analytics here.</div>', unsafe_allow_html=True)
            if st.button("← Back to Quiz", type="secondary"):
                st.session_state.nav_page = "home"; st.rerun()

        # ══════════════════════════════════════════════════════════════════════
        # HISTORY PAGE
        # ══════════════════════════════════════════════════════════════════════
        elif nav == "history":
            hist = st.session_state.quiz_history
            st.markdown("""
            <div class="screen-head">
              <span class="eyebrow">Past Sessions</span>
              <h2>Quiz History</h2>
              <p>All quiz attempts from your current session.</p>
            </div>
            """, unsafe_allow_html=True)
            if not hist:
                st.markdown('<div style="color:rgba(255,255,255,.3);font-size:13px;text-align:center;padding:20px">No quizzes completed yet.</div>', unsafe_allow_html=True)
            else:
                for h in reversed(hist):
                    color = "#34d399" if h["pct"]>=70 else "#fb923c" if h["pct"]>=40 else "#f87171"
                    st.markdown(f"""
                    <div class="hist-row">
                      <div>
                        <div class="hr-title">📄 {h["file"]}</div>
                        <div class="hr-sub">{h["difficulty"]} · {h["total"]} questions · {h["time"]}</div>
                      </div>
                      <div class="hist-score" style="color:{color}">{h["pct"]}%</div>
                    </div>
                    """, unsafe_allow_html=True)
            if st.button("← Back to Quiz", type="secondary"):
                st.session_state.nav_page = "home"; st.rerun()

        # ══════════════════════════════════════════════════════════════════════
        # SETTINGS PAGE
        # ══════════════════════════════════════════════════════════════════════
        elif nav == "settings":
            st.markdown("""
            <div class="screen-head">
              <span class="eyebrow">Configuration</span>
              <h2>Settings</h2>
              <p>Your current app configuration.</p>
            </div>
            """, unsafe_allow_html=True)
            st.markdown(f"""
            <div class="setting-row"><div><div class="sr-label">AI Model</div><div class="sr-sub">Language model used for generation</div></div><span class="sr-val">claude-haiku-4-5</span></div>
            <div class="setting-row"><div><div class="sr-label">Provider</div><div class="sr-sub">API routing service</div></div><span class="sr-val">OpenRouter</span></div>
            <div class="setting-row"><div><div class="sr-label">Max Questions</div><div class="sr-sub">Maximum per quiz session</div></div><span class="sr-val">30</span></div>
            <div class="setting-row"><div><div class="sr-label">File Types</div><div class="sr-sub">Supported upload formats</div></div><span class="sr-val">.pptx</span></div>
            <div class="setting-row"><div><div class="sr-label">Max File Size</div><div class="sr-sub">Upload limit</div></div><span class="sr-val">25 MB</span></div>
            <div class="setting-row"><div><div class="sr-label">AI Explanations</div><div class="sr-sub">Generated for wrong answers</div></div><span class="sr-val">Enabled</span></div>
            """, unsafe_allow_html=True)
            st.write("")
            if st.button("← Back to Quiz", type="secondary"):
                st.session_state.nav_page = "home"; st.rerun()

        # ══════════════════════════════════════════════════════════════════════
        # HOME — QUIZ FLOW
        # ══════════════════════════════════════════════════════════════════════
        elif nav == "home" and st.session_state.step == "upload":
            st.markdown("""
            <div class="glass-panel">
              <div class="screen-head">
                <span class="eyebrow">Step 1 of 3</span>
                <h2>Upload Your Presentation</h2>
                <p>Import a .pptx file — AI will extract every slide, parse all text, and generate smart questions from your content.</p>
              </div>
              <div class="upload-zone">
                <div class="upload-icon-wrap">📤</div>
                <h3>Drag & drop your .pptx here</h3>
                <p>Drop your file anywhere on this card, or use the button below</p>
                <span class="constraint">📎 Supported: .pptx &nbsp;·&nbsp; Max size: 25 MB</span>
              </div>
            </div>
            """, unsafe_allow_html=True)

            up = st.file_uploader("Upload PowerPoint", type=["pptx"], label_visibility="collapsed")

            if up:
                if up.size > 25*1024*1024:
                    st.error("File exceeds 25 MB.")
                else:
                    # Show AI processing animation
                    proc_placeholder = st.empty()
                    stages = ["Extracting slide content…","Analyzing text with AI…","Generating smart questions…"]
                    for stage_idx, stage_txt in enumerate(stages):
                        done_stages = "".join([f'<div class="proc-step done"><div class="step-dot"></div>{stages[j]}</div>' for j in range(stage_idx)])
                        active_stage = f'<div class="proc-step active"><div class="step-dot"></div>{stage_txt}</div>'
                        pending = "".join([f'<div class="proc-step"><div class="step-dot"></div>{stages[j]}</div>' for j in range(stage_idx+1, len(stages))])
                        proc_placeholder.markdown(f"""
                        <div class="processing-wrap">
                          <div class="pulse-ring-wrap">
                            <div class="pulse-ring"></div>
                            <div class="pulse-ring"></div>
                            <div class="pulse-ring"></div>
                            <div class="pulse-core">🧠</div>
                          </div>
                          <div class="processing-title">AI is working…</div>
                          <div class="processing-sub">Please wait while we process your file</div>
                          <div class="proc-steps">{done_stages}{active_stage}{pending}</div>
                        </div>
                        """, unsafe_allow_html=True)
                        time.sleep(0.8)

                    try:
                        slides, wc = extract_pptx(up.read())
                        proc_placeholder.empty()
                        if not slides:
                            st.error("No readable text found in this file.")
                        else:
                            st.session_state.update(slides=slides, filename=up.name, slide_count=len(slides), word_count=wc)
                            st.markdown(f"""
                            <div class="file-card">
                              <span class="fc-icon">📄</span>
                              <div class="fc-info">
                                <strong>{up.name}</strong>
                                <span>{len(slides)} slides extracted &nbsp;·&nbsp; {wc:,} words ready</span>
                              </div>
                              <span class="fc-badge">✓ READY</span>
                            </div>
                            """, unsafe_allow_html=True)
                    except Exception as e:
                        proc_placeholder.empty()
                        st.error(f"Failed to parse: {e}")

            st.write("")
            if st.button("Continue to Configure →", type="primary", disabled=not bool(st.session_state.slides)):
                st.session_state.step = "config"; st.rerun()

        # ══════════════════════════════════════════════════════════════════════
        # SCREEN 2 — CONFIGURE
        # ══════════════════════════════════════════════════════════════════════
        elif nav == "home" and st.session_state.step == "config":
            st.markdown("""
            <div class="screen-head">
              <span class="eyebrow">Step 2 of 3</span>
              <h2>Configure Your Quiz</h2>
              <p>Set difficulty and choose how many questions to generate.</p>
            </div>
            """, unsafe_allow_html=True)

            st.markdown(f'<div class="src-pill">📎 <strong style="color:#f1f5f9">{st.session_state.filename}</strong> <em>&nbsp;·&nbsp; {st.session_state.slide_count} slides &nbsp;·&nbsp; {st.session_state.word_count:,} words</em></div>', unsafe_allow_html=True)

            st.session_state.num_questions = st.slider("Number of questions", 5, 30, st.session_state.num_questions)

            st.markdown('<span class="slbl">Difficulty Level</span>', unsafe_allow_html=True)
            c1,c2,c3 = st.columns(3)
            for col, (name,emoji) in zip([c1,c2,c3],[("Simple","🟢"),("Medium","🟡"),("Complex","🔴")]):
                with col:
                    on = "on" if st.session_state.difficulty==name else ""
                    st.markdown(f'<div class="diff-tile {on}"><span class="de">{emoji}</span><span class="dn">{name}</span></div>', unsafe_allow_html=True)
                    if st.button(name, key=f"d_{name}", type="primary" if on else "secondary"):
                        st.session_state.difficulty=name; st.rerun()

            st.markdown(f'<div class="hint-card" style="margin-top:10px"><span>💡</span><span>{DIFF_HINTS[st.session_state.difficulty]}</span></div>', unsafe_allow_html=True)
            st.write("")

            if st.button(f"⚡  Generate {st.session_state.num_questions} Questions", type="primary"):
                with st.spinner("AI is crafting your quiz — ~15 seconds…"):
                    try:
                        qs = gen_mcqs(st.session_state.slides, st.session_state.num_questions, st.session_state.difficulty)
                        st.session_state.update(questions=qs, user_answers={}, current_q=0, start_time=time.time(), step="quiz", rk=st.session_state.rk+1)
                        st.rerun()
                    except json.JSONDecodeError: st.error("AI returned malformed JSON — retry.")
                    except Exception as e: st.error(f"Error: {e}")

        # ══════════════════════════════════════════════════════════════════════
        # SCREEN 3 — QUIZ (one question per screen)
        # ══════════════════════════════════════════════════════════════════════
        elif nav == "home" and st.session_state.step == "quiz":
            qs      = st.session_state.questions
            total   = len(qs)
            i       = st.session_state.current_q
            q       = qs[i]
            qid     = str(q["id"])
            elapsed = int(time.time()-(st.session_state.start_time or time.time()))
            m,s     = divmod(elapsed,60)
            pct     = int((i/total)*100)
            sel     = st.session_state.user_answers.get(qid)

            st.markdown(f"""
            <div class="quiz-meta">
              <span class="qm-counter">QUESTION {i+1} OF {total}</span>
              <span class="qm-timer">⏱ {m:02d}:{s:02d}</span>
            </div>
            <div class="pbar"><div class="pbar-f" style="width:{pct}%"></div></div>
            <span class="q-topic">{q.get('topic','')}</span>
            <div class="q-text">{q['question']}</div>
            """, unsafe_allow_html=True)

            for key,val in q["options"].items():
                picked = sel == key
                if st.button(f"{key}   {val}", key=f"o_{i}_{key}_{st.session_state.rk}", use_container_width=True, type="primary" if picked else "secondary"):
                    st.session_state.user_answers[qid]=key; st.rerun()

            st.write("")
            cp,_,cn = st.columns([1,.15,1])
            with cp:
                if i>0:
                    if st.button("← Back", type="secondary"):
                        st.session_state.current_q-=1; st.rerun()
            with cn:
                lbl = "Submit & See Results ✓" if i==total-1 else "Next Question →"
                if st.button(lbl, type="primary"):
                    if i<total-1:
                        st.session_state.current_q+=1; st.rerun()
                    else:
                        st.session_state.elapsed=int(time.time()-st.session_state.start_time)
                        with st.spinner("Generating AI explanations…"):
                            try: st.session_state.explanations=gen_exp(qs,st.session_state.user_answers)
                            except: st.session_state.explanations={}
                        st.session_state.step="results"; st.rerun()

        # ══════════════════════════════════════════════════════════════════════
        # SCREEN 4 — RESULTS
        # ══════════════════════════════════════════════════════════════════════
        elif nav == "home" and st.session_state.step == "results":
            qs      = st.session_state.questions
            total   = len(qs)
            correct = sum(1 for q in qs if st.session_state.user_answers.get(str(q["id"]))==q["correct"])
            wrong   = total-correct
            pct     = round((correct/total)*100) if total else 0
            m,s     = divmod(st.session_state.elapsed,60)
            hl      = ("Perfect score! 🎉" if pct==100 else "Excellent! 🌟" if pct>=90 else "Great job! 💪" if pct>=80 else "Good effort 📚" if pct>=60 else "Keep going! 🚀")

            # save to history once
            hist_key = f"{st.session_state.filename}_{st.session_state.elapsed}_{correct}"
            if not any(h.get("_key")==hist_key for h in st.session_state.quiz_history):
                st.session_state.quiz_history.append({
                    "_key": hist_key,
                    "file": st.session_state.filename,
                    "difficulty": st.session_state.difficulty,
                    "total": total, "correct": correct, "pct": pct,
                    "time": f"{m:02d}:{s:02d}",
                })

            st.markdown(f"""
            <div class="result-banner">
              <div class="rb-score">{correct}/{total}</div>
              <div class="rb-pct">{pct}% Score</div>
              <div class="rb-msg">{hl}</div>
              <div class="rb-chips">
                <span class="rbc rbc-g">✓ {correct} Correct</span>
                <span class="rbc rbc-r">✗ {wrong} Wrong</span>
                <span class="rbc rbc-b">⏱ {m:02d}:{s:02d}</span>
              </div>
            </div>
            <div class="rev-section">Review &amp; AI Feedback</div>
            <div class="rev-scroll">
            """, unsafe_allow_html=True)

            for idx,q in enumerate(qs):
                qid=str(q["id"]); ua=st.session_state.user_answers.get(qid)
                ok=ua==q["correct"]; exp=st.session_state.explanations.get(qid,"")
                if ok:   ans=f'<div class="rv-ok">✓ {ua}: {q["options"].get(ua,"")} — correct</div>'
                elif ua: ans=f'<div class="rv-bad">✗ You: {ua}: {q["options"].get(ua,"")} · Correct: {q["correct"]}: {q["options"][q["correct"]]}</div>'
                else:    ans=f'<div class="rv-bad">✗ Not answered · Correct: {q["correct"]}: {q["options"][q["correct"]]}</div>'
                ai=f'<div class="rv-ai"><span>🤖</span><span>{exp}</span></div>' if (not ok and exp) else ""
                st.markdown(f'<div class="rv {"ok" if ok else "bad"}"><div class="rv-hd"><span>{"✅" if ok else "❌"}</span><span class="rv-tl">Q{idx+1} · {q.get("topic","")}</span></div><div class="rv-q">{q["question"]}</div>{ans}{ai}</div>', unsafe_allow_html=True)

            st.markdown('</div>', unsafe_allow_html=True)
            st.write("")
            c1,c2=st.columns(2)
            with c1:
                if st.button("↺  Retake Quiz", type="primary"):
                    st.session_state.update(user_answers={},current_q=0,start_time=time.time(),explanations={},step="quiz",rk=st.session_state.rk+1)
                    st.rerun()
            with c2:
                if st.button("↑  Upload New PPT", type="secondary"):
                    keep = {"quiz_history": st.session_state.quiz_history}
                    for k in list(st.session_state.keys()): del st.session_state[k]
                    st.session_state.quiz_history = keep["quiz_history"]
                    st.rerun()
